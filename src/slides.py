"""Slide rendering: the DETERMINISTIC half of the slide ability.

The model only ever produces a small, forgiving STRUCTURE (title + slides + bullets).
Everything that decides quality - theme, typography, spacing, animation, navigation -
lives here in code, so the deck looks good no matter how weak the model is.

Outputs:
  * a self-contained animated HTML deck (no external assets, works offline)
  * an editable .pptx (python-pptx)
"""
import ast
import json
import re
from datetime import datetime
from pathlib import Path

MAX_SLIDES = 12
MAX_BULLETS = 6


# ---------- tolerant spec parsing (weak models emit sloppy JSON) ----------
def parse_deck(raw):
    """Coerce whatever the model produced into a valid deck spec, or return None.
    Accepts fenced code blocks, leading prose, single quotes, trailing commas."""
    if isinstance(raw, dict):
        return _normalise(raw)
    text = (raw or "").strip()
    fence = re.search(r"```(?:json)?\s*(.+?)```", text, re.S)
    if fence:
        text = fence.group(1).strip()
    start, end = text.find("{"), text.rfind("}")
    if start == -1 or end <= start:
        return None
    blob = text[start:end + 1]
    no_commas = re.sub(r",\s*([}\]])", r"\1", blob)          # trailing commas
    swapped = no_commas.replace("'", '"')                    # single -> double quotes
    for attempt in (blob, no_commas, swapped):
        try:
            return _normalise(json.loads(attempt))
        except Exception:
            continue
    for attempt in (blob, no_commas):                        # python-dict-ish output
        try:
            return _normalise(ast.literal_eval(attempt))     # literals only - no code runs
        except Exception:
            continue
    return None


def _normalise(d):
    """Validate + clamp. Returns None if it isn't usable as a deck."""
    if not isinstance(d, dict):
        return None
    slides_raw = d.get("slides") or d.get("Slides") or []
    if not isinstance(slides_raw, list) or not slides_raw:
        return None
    slides, seen_titles = [], set()
    for s in slides_raw[:MAX_SLIDES]:
        if not isinstance(s, dict):
            continue
        title = str(s.get("title") or s.get("heading") or "").strip()
        key = re.sub(r"\W+", "", title.lower())
        if key and key in seen_titles:        # weak models repeat slides - drop duplicates
            continue
        seen_titles.add(key)
        bl = s.get("bullets") or s.get("points") or s.get("content") or []
        if isinstance(bl, str):
            bl = [x.strip("-*• ").strip() for x in bl.split("\n")]
        bullets = [str(b).strip() for b in bl if str(b).strip()][:MAX_BULLETS]
        if title or bullets:
            slides.append({"title": title or "•", "bullets": bullets,
                           "note": str(s.get("note") or "").strip()})
    if not slides:
        return None
    return {
        "title": str(d.get("title") or "Untitled deck").strip(),
        "subtitle": str(d.get("subtitle") or "").strip(),
        "slides": slides,
    }


def deck_problems(deck):
    """Human-readable validation errors, fed back to the model for the repair loop."""
    if not deck:
        return ["Output was not valid JSON with a non-empty 'slides' list."]
    problems = []
    if len(deck["slides"]) < 3:
        problems.append(f"Only {len(deck['slides'])} slides; produce at least 3.")
    for i, s in enumerate(deck["slides"], 1):
        if not s["bullets"]:
            problems.append(f"Slide {i} ('{s['title'][:30]}') has no bullets.")
    return problems


# ---------- animated, self-contained HTML deck ----------
_HTML = """<!doctype html>
<html lang="en"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>%%TITLE%%</title>
<style>
*{box-sizing:border-box;margin:0;padding:0}
:root{--bg:#0c0c0f;--fg:#f2f2f4;--dim:#a6a6ad;--accent:#e0906a;--accent2:#8b7bff;
      --glass:rgba(255,255,255,.045);--stroke:rgba(255,255,255,.10)}
html,body{height:100%;background:var(--bg);color:var(--fg);
  font-family:'Segoe UI',Inter,system-ui,-apple-system,sans-serif;overflow:hidden}
body::before{content:"";position:fixed;inset:-30%;z-index:0;
  background:radial-gradient(circle at 22% 18%,rgba(224,144,106,.16),transparent 42%),
             radial-gradient(circle at 80% 80%,rgba(139,123,255,.14),transparent 45%);
  animation:drift 22s ease-in-out infinite alternate}
@keyframes drift{to{transform:translate3d(3%,-3%,0) scale(1.08)}}
.deck{position:relative;z-index:1;height:100%}
.slide{position:absolute;inset:0;display:none;flex-direction:column;justify-content:center;
  padding:7vh 9vw;opacity:0}
.slide.active{display:flex;animation:slidein .55s cubic-bezier(.22,.9,.28,1) forwards}
@keyframes slidein{from{opacity:0;transform:translateY(26px) scale(.985)}
                   to{opacity:1;transform:none}}
.slide.cover{align-items:flex-start;justify-content:center}
.eyebrow{font-size:.78rem;letter-spacing:.22em;text-transform:uppercase;color:var(--accent);
  font-weight:600;margin-bottom:1.1rem;opacity:0;animation:rise .5s .05s forwards}
h1{font-size:clamp(2.4rem,6vw,4.6rem);line-height:1.04;font-weight:700;letter-spacing:-.02em;
  background:linear-gradient(120deg,#fff 20%,var(--accent) 95%);-webkit-background-clip:text;
  background-clip:text;color:transparent;opacity:0;animation:rise .6s .12s forwards}
.sub{margin-top:1.4rem;font-size:clamp(1rem,1.8vw,1.35rem);color:var(--dim);max-width:46ch;
  opacity:0;animation:rise .6s .26s forwards}
h2{font-size:clamp(1.7rem,3.6vw,2.9rem);font-weight:650;letter-spacing:-.015em;
  margin-bottom:2.2rem;opacity:0;animation:rise .5s .06s forwards}
h2::after{content:"";display:block;width:74px;height:3px;margin-top:1rem;border-radius:3px;
  background:linear-gradient(90deg,var(--accent),var(--accent2))}
ul{list-style:none;display:flex;flex-direction:column;gap:1.05rem;max-width:62ch}
li{position:relative;padding-left:2rem;font-size:clamp(1rem,1.75vw,1.4rem);line-height:1.5;
  color:#e6e6ea;opacity:0;animation:rise .5s forwards}
li::before{content:"";position:absolute;left:0;top:.62em;width:9px;height:9px;border-radius:50%;
  background:linear-gradient(135deg,var(--accent),var(--accent2));
  box-shadow:0 0 14px rgba(224,144,106,.6)}
@keyframes rise{from{opacity:0;transform:translateY(16px)}to{opacity:1;transform:none}}
.bar{position:fixed;left:0;bottom:0;height:3px;z-index:3;
  background:linear-gradient(90deg,var(--accent),var(--accent2));transition:width .45s ease}
.hud{position:fixed;right:22px;bottom:16px;z-index:3;font-size:.78rem;color:var(--dim);
  background:var(--glass);border:1px solid var(--stroke);border-radius:999px;padding:5px 13px;
  backdrop-filter:blur(12px)}
.hint{position:fixed;left:22px;bottom:16px;z-index:3;font-size:.72rem;color:#6d6d75}
</style></head><body>
<div class="deck" id="deck">%%SLIDES%%</div>
<div class="bar" id="bar"></div><div class="hud" id="hud"></div>
<div class="hint">← → / space · F fullscreen</div>
<script>
const slides=[...document.querySelectorAll('.slide')];let i=0;
function show(n){
  i=Math.max(0,Math.min(slides.length-1,n));
  slides.forEach((s,k)=>s.classList.toggle('active',k===i));
  const li=slides[i].querySelectorAll('li');
  li.forEach((el,k)=>{el.style.animation='none';el.offsetHeight;
    el.style.animation=`rise .5s ${0.18+k*0.10}s forwards`;});
  document.getElementById('bar').style.width=((i+1)/slides.length*100)+'%';
  document.getElementById('hud').textContent=(i+1)+' / '+slides.length;
}
addEventListener('keydown',e=>{
  if(['ArrowRight',' ','PageDown','Enter'].includes(e.key)){e.preventDefault();show(i+1)}
  else if(['ArrowLeft','PageUp'].includes(e.key))show(i-1);
  else if(e.key==='Home')show(0); else if(e.key==='End')show(slides.length-1);
  else if(e.key.toLowerCase()==='f'){document.fullscreenElement?document.exitFullscreen():document.documentElement.requestFullscreen()}
});
addEventListener('click',e=>show(i+(e.clientX < innerWidth*0.22 ? -1 : 1)));
show(0);
</script></body></html>"""


def _esc(s):
    return (str(s).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))


def render_html(deck, path):
    parts = []
    cover_sub = f'<p class="sub">{_esc(deck["subtitle"])}</p>' if deck.get("subtitle") else ""
    parts.append(
        '<section class="slide cover">'
        f'<div class="eyebrow">{datetime.now():%d %b %Y}</div>'
        f'<h1>{_esc(deck["title"])}</h1>{cover_sub}</section>'
    )
    for s in deck["slides"]:
        lis = "".join(f"<li>{_esc(b)}</li>" for b in s["bullets"])
        parts.append(f'<section class="slide"><h2>{_esc(s["title"])}</h2><ul>{lis}</ul></section>')
    html = _HTML.replace("%%TITLE%%", _esc(deck["title"])).replace("%%SLIDES%%", "".join(parts))
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(html, encoding="utf-8")
    return path


# ---------- editable .pptx ----------
def render_pptx(deck, path):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    BG, FG, DIM, ACCENT = RGBColor(0x0C, 0x0C, 0x0F), RGBColor(0xF2, 0xF2, 0xF4), RGBColor(0xA6, 0xA6, 0xAD), RGBColor(0xE0, 0x90, 0x6A)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    def paint(slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG

    cover = prs.slides.add_slide(prs.slide_layouts[6])
    paint(cover)
    tb = cover.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.5), Inches(2.2)).text_frame
    tb.word_wrap = True
    r = tb.paragraphs[0].add_run()
    r.text = deck["title"]
    r.font.size, r.font.bold, r.font.color.rgb, r.font.name = Pt(48), True, FG, "Segoe UI"
    if deck.get("subtitle"):
        p = tb.add_paragraph()
        r = p.add_run()
        r.text = deck["subtitle"]
        r.font.size, r.font.color.rgb, r.font.name = Pt(20), DIM, "Segoe UI"

    for s in deck["slides"]:
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        paint(sl)
        h = sl.shapes.add_textbox(Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.1)).text_frame
        h.word_wrap = True
        r = h.paragraphs[0].add_run()
        r.text = s["title"]
        r.font.size, r.font.bold, r.font.color.rgb, r.font.name = Pt(32), True, ACCENT, "Segoe UI"
        body = sl.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.6)).text_frame
        body.word_wrap = True
        for j, b in enumerate(s["bullets"]):
            p = body.paragraphs[0] if j == 0 else body.add_paragraph()
            p.space_after = Pt(14)
            r = p.add_run()
            r.text = "•  " + b
            r.font.size, r.font.color.rgb, r.font.name = Pt(19), FG, "Segoe UI"
        if s.get("note"):
            sl.notes_slide.notes_text_frame.text = s["note"]

    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path
